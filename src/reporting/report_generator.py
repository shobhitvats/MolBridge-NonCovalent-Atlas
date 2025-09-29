"""
Report generation for Protein Interaction Explorer.
Generates PDF, PowerPoint, and other export formats.
"""

import io
import json
from datetime import datetime
from typing import Dict, Any, List, Optional
from pathlib import Path
import base64

import pandas as pd
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

from utils.config import AppConfig

# Optional extension compute functions (lazy import for packaging). These are lightweight wrappers
# that each expose a compute(result, config) -> dict contract. We only import them here to honor
# the force_package_extensions flag so that exports contain structural/advanced analytics even if
# the UI never visited the lazy-computation tabs.
try:  # Guard so report generation still works if some experimental modules are absent
    from analysis.extensions import (
        compute_residue_profiles,
        compute_interface_analysis,
        compute_outliers,
        compute_provenance,
        compute_motifs,
        compute_secondary_structure,
        compute_sasa_bsa,
        compute_geometry_quality,
        compute_disulfides,
        compute_pockets,
        compute_conservation,
        compute_pi_pi_refinement,
    compute_hbond_subtypes,
    )
except Exception:  # pragma: no cover - defensive
    compute_residue_profiles = compute_interface_analysis = compute_outliers = None
    compute_provenance = compute_motifs = None
    compute_secondary_structure = compute_sasa_bsa = None
    compute_geometry_quality = compute_disulfides = None
    compute_pockets = compute_conservation = compute_pi_pi_refinement = compute_hbond_subtypes = None

class ReportGenerator:
    """Generates comprehensive reports from interaction analysis results."""
    
    def __init__(self, config: AppConfig):
        self.config = config

    # ------------------------------------------------------------------
    # Data Extraction & Advanced Statistics Helpers
    # ------------------------------------------------------------------
    def _iter_interaction_items(self, analysis_results: Dict[str, Dict[str, Any]]):
        """Yield (pdb_id, interaction_type, interaction_obj) triples for all interactions.

        Handles both dict-based and object-based interaction representations.
        """
        for pdb_id, data in analysis_results.items():
            if not data:
                continue
            interactions = data.get('interactions', {})
            for int_type, int_list in interactions.items():
                for obj in int_list:
                    yield pdb_id, int_type, obj

    def _safe_float(self, value: Any) -> Optional[float]:
        try:
            if value is None:
                return None
            f = float(value)
            if f != f:  # NaN guard
                return None
            return f
        except Exception:
            return None

    def _compute_basic_stats(self, values: List[float]) -> Dict[str, Any]:
        if not values:
            return {}
        import math
        n = len(values)
        mean = sum(values) / n
        var = sum((v - mean) ** 2 for v in values) / n if n > 1 else 0.0
        std = math.sqrt(var)
        return {
            'count': n,
            'mean': round(mean, 3),
            'std': round(std, 3),
            'min': round(min(values), 3),
            'max': round(max(values), 3)
        }

    def _compute_interaction_statistics(self, analysis_results: Dict[str, Dict[str, Any]]) -> Dict[str, Any]:
        """Compute advanced statistics across all structures.

        Returns dict with per interaction-type metrics: counts, distance stats, angle stats, strength distribution.
        """
        per_type: Dict[str, Dict[str, Any]] = {}
        global_totals = 0
        for pdb_id, int_type, obj in self._iter_interaction_items(analysis_results):
            per_type.setdefault(int_type, {
                'count': 0,
                'distances': [],
                'angles': [],
                'strengths': {},
                'structures': set()
            })
            rec = per_type[int_type]
            rec['count'] += 1
            global_totals += 1
            rec['structures'].add(pdb_id)
            # Distance / angle extraction
            d = self._safe_float(self._get_interaction_property(obj, 'distance'))
            if d is not None:
                rec['distances'].append(d)
            a = self._safe_float(self._get_interaction_property(obj, 'angle'))
            if a is not None:
                rec['angles'].append(a)
            s = self._get_interaction_property(obj, 'strength')
            if s:
                rec['strengths'][s] = rec['strengths'].get(s, 0) + 1

        # Finalize stats
        finalized = {}
        for int_type, rec in per_type.items():
            finalized[int_type] = {
                'total_count': rec['count'],
                'structures_with': len(rec['structures']),
                'distance_stats': self._compute_basic_stats(rec['distances']),
                'angle_stats': self._compute_basic_stats(rec['angles']),
                'strength_distribution': rec['strengths']
            }

        return {
            'global_total_interactions': global_totals,
            'interaction_types': finalized
        }

    def _build_residue_network(self, analysis_results: Dict[str, Dict[str, Any]]) -> Dict[str, Any]:
        """Construct a residue interaction network (RIN) summary.

        Nodes: residue (ResidueID|Chain). Edges weighted by number of interactions between residue pair (all types).
        Returns top hubs by degree and weighted degree.
        """
        # Edge map: (res1, res2) -> weight
        edges: Dict[tuple, int] = {}
        node_partners: Dict[str, set] = {}
        node_weights: Dict[str, int] = {}

        def norm_res(residue: Any, chain: Any) -> str:
            return f"{residue}_{chain}" if chain else residue

        for _pdb_id, int_type, obj in self._iter_interaction_items(analysis_results):
            r1 = self._get_interaction_property(obj, 'residue1') or self._get_interaction_property(obj, 'donor_residue') or self._get_interaction_property(obj, 'ring1_residue')
            r2 = self._get_interaction_property(obj, 'residue2') or self._get_interaction_property(obj, 'acceptor_residue') or self._get_interaction_property(obj, 'ring2_residue')
            c1 = self._get_interaction_property(obj, 'chain1') or self._get_interaction_property(obj, 'donor_chain') or self._get_interaction_property(obj, 'ring1_chain')
            c2 = self._get_interaction_property(obj, 'chain2') or self._get_interaction_property(obj, 'acceptor_chain') or self._get_interaction_property(obj, 'ring2_chain')
            if not (r1 and r2):
                continue
            n1 = norm_res(r1, c1)
            n2 = norm_res(r2, c2)
            if n1 == n2:
                continue
            key = tuple(sorted((n1, n2)))
            edges[key] = edges.get(key, 0) + 1
        
        for (n1, n2), w in edges.items():
            node_partners.setdefault(n1, set()).add(n2)
            node_partners.setdefault(n2, set()).add(n1)
            node_weights[n1] = node_weights.get(n1, 0) + w
            node_weights[n2] = node_weights.get(n2, 0) + w

        hubs = []
        for node, partners in node_partners.items():
            hubs.append({
                'residue': node,
                'degree': len(partners),
                'weighted_degree': node_weights.get(node, 0)
            })

        # Top hubs by degree then weighted degree
        hubs_sorted = sorted(hubs, key=lambda x: (x['degree'], x['weighted_degree']), reverse=True)[:15]

        network_summary = {
            'nodes': len(node_partners),
            'edges': len(edges),
            'top_hubs': hubs_sorted
        }
        return network_summary
    
    def _get_interaction_property(self, interaction: Any, property_name: str, default_value: Any = None) -> Any:
        """
        Safely get a property from an interaction object, whether it's a dictionary or an object with attributes.

        Args:
            interaction: The interaction object (dict or custom object)
            property_name: Name of the property to retrieve
            default_value: Default value if property is not found

        Returns:
            The property value or default_value
        """
        if hasattr(interaction, property_name):
            # Object-style interaction (e.g., HydrogenBond dataclass)
            return getattr(interaction, property_name, default_value)
        elif isinstance(interaction, dict):
            # Dictionary-style interaction
            return interaction.get(property_name, default_value)
        else:
            # Unknown format, return default
            return default_value
    
    def generate_pdf_report(self,
                           pdb_ids: List[str],
                           analysis_results: Dict[str, Dict[str, Any]],
                           options: Dict[str, Any]) -> bytes:
        """
        Generate comprehensive PDF report.
        
        Args:
            pdb_ids: List of PDB identifiers
            analysis_results: Analysis results for each PDB
            options: Report generation options
            
        Returns:
            PDF report as bytes
        """
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            spaceAfter=12,
            spaceBefore=20
        )
        
        # Cover page
        story.extend(self._create_cover_page(pdb_ids, title_style, styles))
        story.append(PageBreak())
        
        # Table of contents
        if options.get('include_metadata', True):
            story.extend(self._create_table_of_contents(pdb_ids, analysis_results, heading_style))
            story.append(PageBreak())
        
        # Executive summary
        story.extend(self._create_executive_summary(pdb_ids, analysis_results, heading_style, styles))
        story.append(PageBreak())

        # Parameter provenance / analysis parameters section
        if options.get('include_provenance', True):
            story.extend(self._create_parameter_provenance_section(analysis_results, heading_style, styles))
            story.append(PageBreak())

        # Detailed interaction statistics section
        if options.get('include_detailed_stats', True):
            story.extend(self._create_detailed_statistics_section(analysis_results, heading_style, styles))
            story.append(PageBreak())

        # Residue interaction network section
        if options.get('include_residue_network', True):
            story.extend(self._create_residue_network_section(analysis_results, heading_style, styles))
            story.append(PageBreak())
        
        # Methodology section
        if options.get('include_methodology', True):
            story.extend(self._create_methodology_section(heading_style, styles))
            story.append(PageBreak())
        
        # Individual structure reports
        for pdb_id in pdb_ids:
            if pdb_id in analysis_results and analysis_results[pdb_id]:
                story.extend(self._create_structure_report(
                    pdb_id, 
                    analysis_results[pdb_id], 
                    heading_style, 
                    styles,
                    options
                ))
                story.append(PageBreak())
        
        # Comparative analysis (if multiple structures)
        if len(pdb_ids) > 1:
            story.extend(self._create_comparative_analysis(
                pdb_ids,
                analysis_results,
                heading_style,
                styles
            ))
            story.append(PageBreak())
        
        # Appendix
        story.extend(self._create_appendix(pdb_ids, analysis_results, heading_style, styles))
        
        # Build PDF
        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()
    
    def _create_cover_page(self, pdb_ids: List[str], title_style, styles) -> List:
        """Create cover page elements."""
        story = []
        
        # Title
        story.append(Paragraph("Protein Interaction Analysis Report", title_style))
        story.append(Spacer(1, 0.5*inch))
        
        # Structure list
        if len(pdb_ids) == 1:
            story.append(Paragraph(f"Structure: {pdb_ids[0]}", styles['Heading2']))
        else:
            story.append(Paragraph(f"Structures Analyzed: {len(pdb_ids)}", styles['Heading2']))
            structure_list = ", ".join(pdb_ids[:10])
            if len(pdb_ids) > 10:
                structure_list += f" and {len(pdb_ids) - 10} more"
            story.append(Paragraph(structure_list, styles['Normal']))
        
        story.append(Spacer(1, 0.5*inch))
        
        # Metadata
        story.append(Paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", styles['Normal']))
        story.append(Paragraph(f"Software: Protein Interaction Explorer v{self.config.version}", styles['Normal']))
        story.append(Spacer(1, 1*inch))
        
        # Summary logo/image placeholder
        story.append(Paragraph("🧬 Comprehensive Noncovalent Interaction Analysis", styles['Heading3']))
        
        return story
    
    def _create_table_of_contents(self, pdb_ids: List[str], analysis_results: Dict, heading_style) -> List:
        """Create table of contents."""
        story = []
        story.append(Paragraph("Table of Contents", heading_style))
        story.append(Spacer(1, 0.3*inch))
        
        toc_items = [
            "Executive Summary",
            "Methodology",
        ]
        
        # Add individual structures
        for pdb_id in pdb_ids:
            toc_items.append(f"Structure Analysis: {pdb_id}")
        
        if len(pdb_ids) > 1:
            toc_items.append("Comparative Analysis")
        
        toc_items.append("Appendix")
        
        for i, item in enumerate(toc_items, 1):
            story.append(Paragraph(f"{i}. {item}", getSampleStyleSheet()['Normal']))
        
        return story
    
    def _create_executive_summary(self, pdb_ids: List[str], analysis_results: Dict, heading_style, styles) -> List:
        """Create executive summary."""
        story = []
        story.append(Paragraph("Executive Summary", heading_style))
        
        # Calculate overall statistics
        total_interactions = 0
        interaction_type_counts = {}
        successful_analyses = 0
        
        for pdb_id in pdb_ids:
            if pdb_id in analysis_results and analysis_results[pdb_id]:
                successful_analyses += 1
                interactions = analysis_results[pdb_id].get('interactions', {})
                
                for int_type, int_list in interactions.items():
                    count = len(int_list)
                    total_interactions += count
                    interaction_type_counts[int_type] = interaction_type_counts.get(int_type, 0) + count
        
        # Summary text
        summary_text = f"""
        This report presents a comprehensive analysis of noncovalent interactions in {len(pdb_ids)} protein structure(s).
        The analysis identified a total of {total_interactions} interactions across {len(interaction_type_counts)} different 
        interaction types. {successful_analyses} of {len(pdb_ids)} structures were successfully analyzed.
        """
        
        story.append(Paragraph(summary_text, styles['Normal']))
        story.append(Spacer(1, 0.3*inch))
        
        # Summary table
        if interaction_type_counts:
            summary_data = [['Interaction Type', 'Total Count', 'Average per Structure']]
            
            for int_type, count in sorted(interaction_type_counts.items(), key=lambda x: x[1], reverse=True):
                display_name = self._get_interaction_display_name(int_type)
                avg_count = count / successful_analyses if successful_analyses > 0 else 0
                summary_data.append([display_name, str(count), f"{avg_count:.1f}"])
            
            table = Table(summary_data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            
            story.append(table)

        # Hydrogen bond subtype aggregate (if extension present in any structure)
        subtype_aggregate = {}
        total_hbonds_all = 0
        for pdb_id in pdb_ids:
            res = analysis_results.get(pdb_id) or {}
            exts = res.get('extensions', {})
            hbext = exts.get('hbond_subtypes', {})
            counts = hbext.get('counts', {})
            total_hbonds_all += hbext.get('total_hbonds', 0)
            for k,v in counts.items():
                subtype_aggregate[k] = subtype_aggregate.get(k,0) + v
        if subtype_aggregate and total_hbonds_all > 0:
            story.append(Spacer(1, 0.25*inch))
            story.append(Paragraph("Hydrogen Bond Subtype Summary", styles['Heading3']))
            sub_data = [["Subtype","Total Count","% of All H-Bonds"]]
            for k,v in sorted(subtype_aggregate.items(), key=lambda x: x[1], reverse=True):
                pct = (v/total_hbonds_all)*100 if total_hbonds_all else 0
                sub_data.append([k.replace('_',' '), str(v), f"{pct:.1f}%"]) 
            sub_table = Table(sub_data)
            sub_table.setStyle(TableStyle([
                ('BACKGROUND',(0,0),(-1,0), colors.HexColor('#566573')),
                ('TEXTCOLOR',(0,0),(-1,0), colors.whitesmoke),
                ('ALIGN',(0,0),(-1,-1),'CENTER'),
                ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
                ('FONTSIZE',(0,0),(-1,0),10),
                ('BACKGROUND',(0,1),(-1,-1), colors.HexColor('#F2F3F4')),
                ('GRID',(0,0),(-1,-1),0.5, colors.grey)
            ]))
            story.append(sub_table)
        
        return story
    
    def _create_methodology_section(self, heading_style, styles) -> List:
        """Create methodology section."""
        story = []
        story.append(Paragraph("Methodology", heading_style))
        
        methodology_text = f"""
        <b>Analysis Parameters:</b><br/>
        This analysis was performed using the Protein Interaction Explorer (version {self.config.version}),
        which implements state-of-the-art algorithms for detecting noncovalent interactions in protein structures.
        
        <br/><br/>
        <b>Interaction Types Analyzed:</b><br/>
        • Hydrogen bonds (conventional, low-barrier, C5-type)<br/>
        • Halogen bonds (Cl, Br, I, F interactions)<br/>
        • π-π stacking interactions<br/>
        • Ionic interactions<br/>
        • Hydrophobic contacts<br/>
        • C-H···π interactions<br/>
        • Chalcogen bonds<br/>
        • Pnictogen bonds<br/>
        • Tetrel bonds<br/>
        • Anion-π interactions<br/>
        • n→π* interactions<br/>
        • London dispersion forces (heuristic)<br/>
        
        <br/>
        <b>Detection Criteria:</b><br/>
        All interactions were detected using literature-based geometric criteria including distance and angle cutoffs.
        The analysis includes both intra- and inter-molecular interactions, with specific handling for different
        assembly types (biological assembly vs. asymmetric unit).
        
        <br/><br/>
        <b>Quality Control:</b><br/>
        Structures underwent validation checks including completeness assessment and clash detection.
        Missing hydrogen atoms were added using standard geometric algorithms where applicable.
        """
        
        story.append(Paragraph(methodology_text, styles['Normal']))
        
        return story

    # ------------------------------------------------------------------
    # New Advanced Sections
    # ------------------------------------------------------------------
    def _create_parameter_provenance_section(self, analysis_results: Dict[str, Dict[str, Any]], heading_style, styles) -> List:
        story = []
        story.append(Paragraph("Parameter & Provenance", heading_style))
        cfg = self.config.interactions
        provenance_html = [
            "<b>Detection Parameter Summary</b><br/>",
            f"Hydrogen Bond Cutoffs: distance ≤ {cfg.hbond_distance_cutoff} Å, angle ≥ {cfg.hbond_angle_cutoff}°<br/>",
            f"Halogen Bond Cutoffs: distance ≤ {cfg.halogen_distance_cutoff} Å, angle ≥ {cfg.halogen_angle_cutoff}°<br/>",
            f"Chalcogen/Pnictogen/Tetrel: distance ≤ 4.0 Å baseline; angle ≥ 140° (specific per type)<br/>",
            f"π-π Stacking: center distance ≤ {cfg.pi_pi_distance_cutoff} Å, plane angle ≤ {cfg.pi_pi_angle_cutoff}°<br/>",
            f"C-H···π: distance {cfg.ch_pi_min_distance}-{cfg.ch_pi_max_distance} Å, angle ≤ {cfg.ch_pi_angle_cutoff}°, ring height ≤ {cfg.ch_pi_max_height} Å<br/>",
            f"Anion-π: centroid distance ≤ {cfg.anion_pi_distance_cutoff} Å<br/>",
            f"n→π*: distance ≤ {cfg.n_pi_star_distance_cutoff} Å, angle ≥ {cfg.n_pi_star_angle_cutoff}°<br/>",
            f"Ionic: charge center distance ≤ {cfg.ionic_distance_cutoff} Å<br/>",
            f"Hydrophobic: heavy atom distance ≤ {cfg.hydrophobic_distance_cutoff} Å (excluding local i,i±1,2 pairs)<br/>",
            "<br/><b>Structure Processing</b><br/>",
            f"Assembly Mode: {self.config.default_assembly.title()}<br/>",
            f"Include Ligands: {self.config.include_ligands}, Exclude Waters: {self.config.exclude_waters}<br/>",
            f"Version: {self.config.version} | Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}<br/>",
        ]
        story.append(Paragraph("".join(provenance_html), styles['Normal']))
        return story

    def _create_detailed_statistics_section(self, analysis_results: Dict[str, Dict[str, Any]], heading_style, styles) -> List:
        story = []
        story.append(Paragraph("Detailed Interaction Statistics", heading_style))
        stats = self._compute_interaction_statistics(analysis_results)
        itypes = stats['interaction_types']
        if not itypes:
            story.append(Paragraph("No interactions available for statistical summary.", styles['Normal']))
            return story

        # Build table header
        table_data = [[
            'Interaction', 'Total', 'Structures',
            'Dist Mean', 'Dist Std', 'Dist Min', 'Dist Max',
            'Ang Mean', 'Ang Std', 'Strength Dist.'
        ]]

        for int_type, rec in sorted(itypes.items(), key=lambda x: x[1]['total_count'], reverse=True):
            disp = self._get_interaction_display_name(int_type)
            d_stats = rec['distance_stats'] or {}
            a_stats = rec['angle_stats'] or {}
            strength_dist = ", ".join(f"{k}:{v}" for k, v in sorted(rec['strength_distribution'].items())) or '—'
            table_data.append([
                disp,
                rec['total_count'],
                rec['structures_with'],
                d_stats.get('mean', '—'),
                d_stats.get('std', '—'),
                d_stats.get('min', '—'),
                d_stats.get('max', '—'),
                a_stats.get('mean', '—'),
                a_stats.get('std', '—'),
                strength_dist
            ])

        table = Table(table_data, repeatRows=1)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#2C3E50')),
            ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
            ('ALIGN', (0,0), (-1,0), 'CENTER'),
            ('FONTSIZE', (0,0), (-1,0), 9),
            ('BACKGROUND', (0,1), (-1,-1), colors.HexColor('#ECF0F1')),
            ('FONTSIZE', (0,1), (-1,-1), 8),
            ('GRID', (0,0), (-1,-1), 0.25, colors.grey),
            ('VALIGN', (0,0), (-1,-1), 'MIDDLE')
        ]))
        story.append(table)
        return story

    def _create_residue_network_section(self, analysis_results: Dict[str, Dict[str, Any]], heading_style, styles) -> List:
        story = []
        story.append(Paragraph("Residue Interaction Network", heading_style))
        network = self._build_residue_network(analysis_results)
        if not network.get('nodes'):
            story.append(Paragraph("Insufficient interaction density to construct a residue network.", styles['Normal']))
            return story
        summary = f"Nodes: {network['nodes']} | Edges: {network['edges']}<br/>Top hubs (residues with highest connectivity):"
        story.append(Paragraph(summary, styles['Normal']))
        hub_data = [['Residue', 'Degree', 'Weighted Degree']]
        for hub in network['top_hubs']:
            hub_data.append([hub['residue'], hub['degree'], hub['weighted_degree']])
        hub_table = Table(hub_data, repeatRows=1)
        hub_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#34495E')),
            ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
            ('ALIGN', (0,0), (-1,-1), 'CENTER'),
            ('FONTSIZE', (0,0), (-1,0), 9),
            ('BACKGROUND', (0,1), (-1,-1), colors.HexColor('#F8F9F9')),
            ('GRID', (0,0), (-1,-1), 0.25, colors.grey)
        ]))
        story.append(Spacer(1, 0.15*inch))
        story.append(hub_table)
        return story
    
    def _create_structure_report(self, 
                               pdb_id: str, 
                               result: Dict[str, Any], 
                               heading_style, 
                               styles,
                               options: Dict[str, Any]) -> List:
        """Create report section for individual structure."""
        story = []
        story.append(Paragraph(f"Structure Analysis: {pdb_id}", heading_style))
        
        # Structure information
        structure_info = result.get('structure_info', {})
        metadata = result.get('metadata', {})
        
        info_text = f"""
        <b>PDB ID:</b> {pdb_id}<br/>
        <b>Analysis Time:</b> {metadata.get('analysis_time', 0):.2f} seconds<br/>
        <b>Total Chains:</b> {len(structure_info.get('chains', []))}<br/>
        <b>Total Residues:</b> {structure_info.get('total_residues', 0)}<br/>
        <b>Total Atoms:</b> {structure_info.get('total_atoms', 0)}<br/>
        """
        
        story.append(Paragraph(info_text, styles['Normal']))
        story.append(Spacer(1, 0.2*inch))
        
        # Interaction summary
        interactions = result.get('interactions', {})
        total_interactions = sum(len(int_list) for int_list in interactions.values())
        
        story.append(Paragraph(f"Total Interactions Detected: {total_interactions}", styles['Heading3']))
        
        # Detailed interaction table
        if total_interactions > 0:
            interaction_data = [['Interaction Type', 'Count', 'Average Distance (Å)']]
            
            for int_type, int_list in interactions.items():
                if int_list:
                    display_name = self._get_interaction_display_name(int_type)
                    count = len(int_list)
                    
                    distances = [self._get_interaction_property(i, 'distance', 0) for i in int_list if self._get_interaction_property(i, 'distance', None) is not None]
                    avg_distance = sum(distances) / len(distances) if distances else 'N/A'
                    avg_distance_str = f"{avg_distance:.2f}" if isinstance(avg_distance, float) else avg_distance
                    
                    interaction_data.append([display_name, str(count), avg_distance_str])
            
            if len(interaction_data) > 1:
                table = Table(interaction_data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.lightblue),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.lightgrey),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                
                story.append(table)
        
        # Hotspots
        hotspots = result.get('hotspots', [])
        if hotspots:
            story.append(Spacer(1, 0.2*inch))
            story.append(Paragraph("Interaction Hotspots", styles['Heading3']))
            
            hotspot_text = "Top interaction hotspots (residues with highest interaction density):<br/>"
            for i, hotspot in enumerate(hotspots[:5], 1):
                hotspot_text += f"{i}. {hotspot['residue']} ({hotspot['interaction_count']} interactions)<br/>"
            
            story.append(Paragraph(hotspot_text, styles['Normal']))

        # Per-structure hydrogen bond subtype table (if extension computed)
        exts = result.get('extensions', {})
        hbext = exts.get('hbond_subtypes') if isinstance(exts, dict) else None
        if hbext and hbext.get('total_hbonds',0) > 0:
            story.append(Spacer(1, 0.2*inch))
            story.append(Paragraph("Hydrogen Bond Subtypes", styles['Heading3']))
            counts = hbext.get('counts', {})
            total_hb = hbext.get('total_hbonds', 0)
            rows = [["Subtype","Count","Fraction"]]
            fractions = hbext.get('fractions', {})
            for k,v in sorted(counts.items(), key=lambda x: x[1], reverse=True):
                rows.append([k.replace('_',' '), str(v), f"{fractions.get(k,0.0):.2f}"])
            table = Table(rows)
            table.setStyle(TableStyle([
                ('BACKGROUND',(0,0),(-1,0), colors.HexColor('#5D6D7E')),
                ('TEXTCOLOR',(0,0),(-1,0), colors.whitesmoke),
                ('ALIGN',(0,0),(-1,-1),'CENTER'),
                ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
                ('FONTSIZE',(0,0),(-1,0),9),
                ('BACKGROUND',(0,1),(-1,-1), colors.HexColor('#EBEDEF')),
                ('GRID',(0,0),(-1,-1),0.5, colors.grey)
            ]))
            story.append(table)
        
        return story
    
    def _create_comparative_analysis(self, 
                                   pdb_ids: List[str],
                                   analysis_results: Dict[str, Dict[str, Any]],
                                   heading_style,
                                   styles) -> List:
        """Create comparative analysis section."""
        story = []
        story.append(Paragraph("Comparative Analysis", heading_style))
        
        # Create comparison table
        comparison_data = [['PDB ID', 'Total Interactions'] + [
            self._get_interaction_display_name(int_type) 
            for int_type in ['hydrogen_bond', 'halogen_bond', 'pi_pi', 'ionic', 'hydrophobic']
        ]]
        
        for pdb_id in pdb_ids:
            if pdb_id in analysis_results and analysis_results[pdb_id]:
                interactions = analysis_results[pdb_id].get('interactions', {})
                total = sum(len(int_list) for int_list in interactions.values())
                
                row = [pdb_id, str(total)]
                for int_type in ['hydrogen_bond', 'halogen_bond', 'pi_pi', 'ionic', 'hydrophobic']:
                    count = len(interactions.get(int_type, []))
                    row.append(str(count))
                
                comparison_data.append(row)
        
        if len(comparison_data) > 1:
            table = Table(comparison_data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.darkblue),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 9),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.lightsteelblue),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            
            story.append(table)
        
        return story
    
    def _create_appendix(self, pdb_ids: List[str], analysis_results: Dict, heading_style, styles) -> List:
        """Create appendix with detailed data."""
        story = []
        story.append(Paragraph("Appendix", heading_style))
        
        appendix_text = f"""
        <b>Software Information:</b><br/>
        Protein Interaction Explorer v{self.config.version}<br/>
        Analysis Date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}<br/>
        
        <br/>
        <b>Data Provenance:</b><br/>
        Structures downloaded from RCSB Protein Data Bank<br/>
        Analysis performed using literature-based geometric criteria<br/>
        
        <br/>
        <b>Citation:</b><br/>
        If you use this analysis in your research, please cite:<br/>
        Protein Interaction Explorer: Comprehensive Analysis of Noncovalent Interactions in Protein Structures<br/>
        """
        
        story.append(Paragraph(appendix_text, styles['Normal']))
        
        return story
    
    def generate_csv_report(self, pdb_id: str, result: Dict[str, Any]) -> str:
        """Generate CSV report for a single structure."""
        interactions = result.get('interactions', {})
        
        # Combine all interactions into a single list
        all_interactions = []
        
        for interaction_type, interaction_list in interactions.items():
            for interaction in interaction_list:
                interaction_data = {
                    'PDB_ID': pdb_id,
                    'Interaction_Type': self._get_interaction_display_name(interaction_type),
                    'Residue_1': self._get_interaction_property(interaction, 'residue1', ''),
                    'Chain_1': self._get_interaction_property(interaction, 'chain1', ''),
                    'Residue_2': self._get_interaction_property(interaction, 'residue2', ''),
                    'Chain_2': self._get_interaction_property(interaction, 'chain2', ''),
                    'Distance_A': self._get_interaction_property(interaction, 'distance', ''),
                    'Angle_deg': self._get_interaction_property(interaction, 'angle', ''),
                    'Strength': self._get_interaction_property(interaction, 'strength', ''),
                    'Atom_1': self._get_interaction_property(interaction, 'atom1', ''),
                    'Atom_2': self._get_interaction_property(interaction, 'atom2', '')
                }
                all_interactions.append(interaction_data)
        
        # Convert to DataFrame and then CSV
        df = pd.DataFrame(all_interactions)
        return df.to_csv(index=False)
    
    def generate_excel_report(self, 
                            pdb_ids: List[str],
                            analysis_results: Dict[str, Dict[str, Any]]) -> bytes:
        """Generate Excel report with multiple sheets."""
        buffer = io.BytesIO()
        
        with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
            # Summary sheet
            summary_data = []
            for pdb_id in pdb_ids:
                if pdb_id in analysis_results and analysis_results[pdb_id]:
                    interactions = analysis_results[pdb_id].get('interactions', {})
                    total = sum(len(int_list) for int_list in interactions.values())
                    
                    row = {'PDB_ID': pdb_id, 'Total_Interactions': total}
                    for int_type, int_list in interactions.items():
                        display_name = self._get_interaction_display_name(int_type).replace(' ', '_')
                        row[display_name] = len(int_list)
                    
                    summary_data.append(row)
            
            if summary_data:
                summary_df = pd.DataFrame(summary_data)
                summary_df.to_excel(writer, sheet_name='Summary', index=False)
            
            # Individual sheets for each structure
            for pdb_id in pdb_ids:
                if pdb_id in analysis_results and analysis_results[pdb_id]:
                    csv_data = self.generate_csv_report(pdb_id, analysis_results[pdb_id])
                    df = pd.read_csv(io.StringIO(csv_data))
                    
                    if not df.empty:
                        sheet_name = f"Interactions_{pdb_id}"[:31]  # Excel sheet name limit
                        df.to_excel(writer, sheet_name=sheet_name, index=False)
        
        buffer.seek(0)
        return buffer.getvalue()
    
    def generate_powerpoint_report(self,
                                 pdb_ids: List[str],
                                 analysis_results: Dict[str, Dict[str, Any]]) -> bytes:
        """Generate PowerPoint presentation."""
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Protein Interaction Analysis"
        subtitle.text = f"Analysis of {len(pdb_ids)} Structure(s)\n{datetime.now().strftime('%Y-%m-%d')}"
        
        # Summary slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = 'Analysis Summary'
        
        tf = body_shape.text_frame
        tf.text = f'Structures Analyzed: {len(pdb_ids)}'
        
        # Calculate totals
        total_interactions = 0
        for pdb_id in pdb_ids:
            if pdb_id in analysis_results and analysis_results[pdb_id]:
                interactions = analysis_results[pdb_id].get('interactions', {})
                total_interactions += sum(len(int_list) for int_list in interactions.values())
        
        p = tf.add_paragraph()
        p.text = f'Total Interactions: {total_interactions}'
        
        # Individual structure slides
        for pdb_id in pdb_ids:
            if pdb_id in analysis_results and analysis_results[pdb_id]:
                self._add_structure_slide(prs, pdb_id, analysis_results[pdb_id])
        
        # Save to buffer
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def generate_complete_package(self,
                                  pdb_ids: List[str],
                                  analysis_results: Dict[str, Dict[str, Any]],
                                  options: Optional[Dict[str, Any]] = None) -> bytes:
        import zipfile, csv
        options = options or {}
        # Force compute extensions if configured
        if getattr(self.config, 'force_package_extensions', False):
            self._ensure_all_extensions(pdb_ids, analysis_results)

        zip_buffer = io.BytesIO()
        stats = self._compute_interaction_statistics(analysis_results)
        network = self._build_residue_network(analysis_results)
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
            # --- Core artifact generation reused from previous implementation (indented into method) ---
            try:
                pdf_bytes = self.generate_pdf_report(pdb_ids, analysis_results, options)
                zf.writestr("report/Protein_Interaction_Report.pdf", pdf_bytes)
            except Exception as e:
                zf.writestr("report/ERROR_pdf.txt", f"Failed to generate PDF: {e}")
            try:
                pptx_bytes = self.generate_powerpoint_report(pdb_ids, analysis_results)
                zf.writestr("report/Protein_Interaction_Presentation.pptx", pptx_bytes)
            except Exception as e:
                zf.writestr("report/ERROR_pptx.txt", f"Failed to generate PPTX: {e}")
            try:
                excel_bytes = self.generate_excel_report(pdb_ids, analysis_results)
                zf.writestr("data/Interaction_Summary.xlsx", excel_bytes)
            except Exception as e:
                zf.writestr("data/ERROR_excel.txt", f"Failed to generate Excel: {e}")
            try:
                latex_text = self.generate_latex_export(pdb_ids, analysis_results)
                zf.writestr("latex/report.tex", latex_text)
            except Exception as e:
                zf.writestr("latex/ERROR_latex.txt", f"Failed to generate LaTeX: {e}")
            try:
                zf.writestr("data/interaction_statistics.json", json.dumps(stats, indent=2))
            except Exception as e:
                zf.writestr("data/ERROR_stats.txt", f"Failed to write statistics: {e}")
            try:
                hub_csv_io = io.StringIO(); writer = csv.writer(hub_csv_io)
                writer.writerow(["Residue","Degree","Weighted_Degree"])
                for hub in network.get('top_hubs', []):
                    writer.writerow([hub['residue'], hub['degree'], hub['weighted_degree']])
                zf.writestr("data/residue_network_hubs.csv", hub_csv_io.getvalue())
            except Exception as e:
                zf.writestr("data/ERROR_network.txt", f"Failed to write network hubs: {e}")
            # Per-structure artifacts (same logic as before)
            for pdb_id in pdb_ids:
                if pdb_id not in analysis_results or not analysis_results[pdb_id]:
                    continue
                try:
                    csv_text = self.generate_csv_report(pdb_id, analysis_results[pdb_id])
                    zf.writestr(f"data/{pdb_id}_interactions.csv", csv_text)
                    summary_json = json.dumps(analysis_results[pdb_id], indent=2, default=str)
                    zf.writestr(f"data/{pdb_id}_summary.json", summary_json)
                    exts = analysis_results[pdb_id].get('extensions', {})
                    # (All extension export blocks preserved)
                    # --- residue_profiles ---
                    if 'residue_profiles' in exts:
                        try:
                            import csv as _csv, io as _io
                            rp_io = _io.StringIO(); writer = _csv.writer(rp_io)
                            writer.writerow(["Residue","Total","TopTypes","Strong","Moderate","Weak"])
                            for residue, rec in exts['residue_profiles'].get('profiles', {}).items():
                                by_type_sorted = sorted(rec.get('by_type', {}).items(), key=lambda x: x[1], reverse=True)[:4]
                                top_types = ';'.join(f"{k}:{v}" for k,v in by_type_sorted)
                                strengths = rec.get('strengths', {})
                                writer.writerow([residue, rec.get('total',0), top_types, strengths.get('strong',0), strengths.get('moderate',0), strengths.get('weak', strengths.get('very_weak',0))])
                            zf.writestr(f"data/{pdb_id}_residue_profiles.csv", rp_io.getvalue())
                        except Exception as ee:
                            zf.writestr(f"data/ERROR_{pdb_id}_residue_profiles.txt", str(ee))
                    # interface_analysis
                    if 'interface_analysis' in exts:
                        try:
                            import csv as _csv, io as _io
                            if_io = _io.StringIO(); writer = _csv.writer(if_io)
                            writer.writerow(["ChainPair","InteractionCount","ApproxInterfaceResidues"])
                            for rec in exts['interface_analysis'].get('interfaces', []):
                                writer.writerow([rec.get('chain_pair'), rec.get('interaction_count'), rec.get('approx_interface_residues')])
                            zf.writestr(f"data/{pdb_id}_interfaces.csv", if_io.getvalue())
                        except Exception as ee:
                            zf.writestr(f"data/ERROR_{pdb_id}_interfaces.txt", str(ee))
                    # outliers
                    if 'outliers' in exts:
                        try:
                            import csv as _csv, io as _io
                            o_io = _io.StringIO(); writer = _csv.writer(o_io)
                            writer.writerow(["Type","Residue1","Residue2","Distance","Angle","Notes"])
                            for rec in exts['outliers'].get('records', []):
                                writer.writerow([rec.get('interaction_type'), rec.get('residue1'), rec.get('residue2'), rec.get('distance'), rec.get('angle'), rec.get('notes')])
                            zf.writestr(f"data/{pdb_id}_outliers.csv", o_io.getvalue())
                        except Exception as ee:
                            zf.writestr(f"data/ERROR_{pdb_id}_outliers.txt", str(ee))
                    # motifs
                    if 'motifs' in exts:
                        try:
                            import csv as _csv, io as _io
                            m_io = _io.StringIO(); writer = _csv.writer(m_io)
                            writer.writerow(["Type","StructureID","CentralResidue","AromaticResidues","Count"])
                            for rec in exts['motifs'].get('motifs', []):
                                writer.writerow([rec.get('type'), rec.get('structure_id'), rec.get('central_residue'), ';'.join(rec.get('aromatic_residues', [])), rec.get('count')])
                            zf.writestr(f"data/{pdb_id}_motifs.csv", m_io.getvalue())
                        except Exception as ee:
                            zf.writestr(f"data/ERROR_{pdb_id}_motifs.txt", str(ee))
                    # Structural/quality exports (secondary_structure, geometry_quality, sasa_bsa, disulfides, pockets, conservation, pi_pi_refinement)
                    # (Reuse logic previously present)
                    if 'secondary_structure' in exts:
                        try:
                            import csv as _csv, io as _io
                            ss = exts['secondary_structure']; ss_io = _io.StringIO(); w = _csv.writer(ss_io)
                            counts = ss.get('counts', {}); fracs = ss.get('fractions', {})
                            w.writerow(["Element","Count","Fraction"])
                            for k in ['H','E','C']:
                                w.writerow([k, counts.get(k,0), fracs.get(f"{k}_frac",0.0)])
                            zf.writestr(f"data/{pdb_id}_secondary_structure_summary.csv", ss_io.getvalue())
                            tors = ss.get('torsions', [])
                            if tors:
                                tors_io = _io.StringIO(); w2 = _csv.writer(tors_io)
                                w2.writerow(["Residue","Chain","Phi","Psi","Omega","SS"])
                                for row in tors[:5000]:
                                    w2.writerow([row.get('residue'), row.get('chain'), row.get('phi'), row.get('psi'), row.get('omega'), row.get('ss')])
                                zf.writestr(f"data/{pdb_id}_torsions.csv", tors_io.getvalue())
                        except Exception as ee:
                            zf.writestr(f"data/ERROR_{pdb_id}_secondary_structure.txt", str(ee))
                    if 'geometry_quality' in exts:
                        try:
                            import csv as _csv, io as _io
                            gq = exts['geometry_quality']; rama = gq.get('ramachandran', {})
                            rq_io = _io.StringIO(); w = _csv.writer(rq_io)
                            w.writerow(["Favored","Allowed","Outliers"]); w.writerow([rama.get('favored',0), rama.get('allowed',0), rama.get('outliers',0)])
                            zf.writestr(f"data/{pdb_id}_ramachandran_summary.csv", rq_io.getvalue())
                            clashes = gq.get('clashes', [])
                            if clashes:
                                cl_io = _io.StringIO(); w2 = _csv.writer(cl_io); header = sorted({k for c in clashes for k in c.keys()})
                                w2.writerow(header)
                                for c in clashes[:5000]: w2.writerow([c.get(h,'') for h in header])
                                zf.writestr(f"data/{pdb_id}_clashes.csv", cl_io.getvalue())
                            tors_out = gq.get('torsion_outliers', [])
                            if tors_out:
                                to_io = _io.StringIO(); w3 = _csv.writer(to_io); header = sorted({k for t in tors_out for k in t.keys()})
                                w3.writerow(header)
                                for t in tors_out: w3.writerow([t.get(h,'') for h in header])
                                zf.writestr(f"data/{pdb_id}_torsion_outliers.csv", to_io.getvalue())
                        except Exception as ee:
                            zf.writestr(f"data/ERROR_{pdb_id}_geometry_quality.txt", str(ee))
                    if 'sasa_bsa' in exts:
                        try:
                            import csv as _csv, io as _io
                            sas = exts['sasa_bsa']; chain_sasa = sas.get('chain_sasa', {})
                            if chain_sasa:
                                cs_io = _io.StringIO(); w = _csv.writer(cs_io); w.writerow(["Chain","SASA"])
                                for ch, rec in chain_sasa.items(): w.writerow([ch, rec.get('sasa',0.0)])
                                zf.writestr(f"data/{pdb_id}_chain_sasa.csv", cs_io.getvalue())
                            bsa = sas.get('buried_surface', {})
                            if bsa:
                                bsa_io = _io.StringIO(); w2 = _csv.writer(bsa_io); w2.writerow(["ChainPair","Buried_SASA"])
                                for pair, val in bsa.items(): w2.writerow([pair, val])
                                zf.writestr(f"data/{pdb_id}_buried_surface.csv", bsa_io.getvalue())
                        except Exception as ee:
                            zf.writestr(f"data/ERROR_{pdb_id}_sasa_bsa.txt", str(ee))
                    if 'disulfides' in exts:
                        try:
                            import csv as _csv, io as _io
                            ds = exts['disulfides'].get('bonds', []); ds_io = _io.StringIO(); w = _csv.writer(ds_io)
                            if ds:
                                header = sorted({k for d in ds for k in d.keys()}); w.writerow(header)
                                for d in ds: w.writerow([d.get(h,'') for h in header])
                            zf.writestr(f"data/{pdb_id}_disulfides.csv", ds_io.getvalue())
                        except Exception as ee:
                            zf.writestr(f"data/ERROR_{pdb_id}_disulfides.txt", str(ee))
                    if 'pockets' in exts:
                        try:
                            import csv as _csv, io as _io
                            pockets = exts['pockets'].get('pockets', []); pk_io = _io.StringIO(); w = _csv.writer(pk_io)
                            if pockets:
                                header = sorted({k for p in pockets for k in p.keys()}); w.writerow(header)
                                for p in pockets[:5000]: w.writerow([p.get(h,'') for h in header])
                            zf.writestr(f"data/{pdb_id}_pockets.csv", pk_io.getvalue())
                        except Exception as ee:
                            zf.writestr(f"data/ERROR_{pdb_id}_pockets.txt", str(ee))
                    if 'conservation' in exts:
                        try:
                            import csv as _csv, io as _io
                            cons = exts['conservation'].get('residues', {}); c_io = _io.StringIO(); w = _csv.writer(c_io)
                            if cons:
                                w.writerow(["Residue","Score","Rank"])
                                for res, rec in cons.items(): w.writerow([res, rec.get('score'), rec.get('rank')])
                                zf.writestr(f"data/{pdb_id}_conservation.csv", c_io.getvalue())
                        except Exception as ee:
                            zf.writestr(f"data/ERROR_{pdb_id}_conservation.txt", str(ee))
                    if 'pi_pi_refinement' in exts:
                        try:
                            import csv as _csv, io as _io
                            ref = exts['pi_pi_refinement']; subtype_io = _io.StringIO(); w = _csv.writer(subtype_io)
                            sub = ref.get('subtypes', {})
                            w.writerow(["Subtype","Count"])
                            for k,v in sub.items(): w.writerow([k,v])
                            zf.writestr(f"data/{pdb_id}_pi_pi_subtypes.csv", subtype_io.getvalue())
                            details = ref.get('details', [])
                            if details:
                                d_io = _io.StringIO(); w2 = _csv.writer(d_io); header = sorted({k for d in details for k in d.keys()})
                                w2.writerow(header)
                                for d in details[:5000]: w2.writerow([d.get(h,'') for h in header])
                                zf.writestr(f"data/{pdb_id}_pi_pi_details.csv", d_io.getvalue())
                        except Exception as ee:
                            zf.writestr(f"data/ERROR_{pdb_id}_pi_pi_refinement.txt", str(ee))
                    if 'hbond_subtypes' in exts:
                        try:
                            import csv as _csv, io as _io
                            hb = exts['hbond_subtypes']; hb_io = _io.StringIO(); w = _csv.writer(hb_io)
                            counts = hb.get('counts', {})
                            w.writerow(["Subtype","Count","Fraction"])
                            fractions = hb.get('fractions', {})
                            total_hb = hb.get('total_hbonds', 0)
                            for k,v in counts.items():
                                w.writerow([k, v, f"{fractions.get(k,0.0):.4f}"])
                            zf.writestr(f"data/{pdb_id}_hbond_subtypes.csv", hb_io.getvalue())
                            annotated = hb.get('annotated', [])
                            if annotated:
                                ann_io = _io.StringIO(); w2 = _csv.writer(ann_io); header = sorted({h for a in annotated for h in a.keys()})
                                w2.writerow(header)
                                for a in annotated[:5000]: w2.writerow([a.get(h,'') for h in header])
                                zf.writestr(f"data/{pdb_id}_hbond_annotated.csv", ann_io.getvalue())
                        except Exception as ee:
                            zf.writestr(f"data/ERROR_{pdb_id}_hbond_subtypes.txt", str(ee))
                except Exception as e:
                    zf.writestr(f"data/ERROR_{pdb_id}.txt", f"Failed to add artifacts for {pdb_id}: {e}")
            # Aggregates
            try:
                if len(pdb_ids) > 1:
                    import csv as _csv, io as _io
                    ss_rows=[]; clash_rows=[]; pocket_rows=[]
                    hbond_rows=[]
                    for pid in pdb_ids:
                        res = analysis_results.get(pid) or {}; exts = res.get('extensions', {})
                        ss = exts.get('secondary_structure', {}); gq = exts.get('geometry_quality', {}); pk = exts.get('pockets', {})
                        hb = exts.get('hbond_subtypes', {})
                        counts = ss.get('counts', {}); fracs = ss.get('fractions', {})
                        ss_rows.append([pid, counts.get('H',0), counts.get('E',0), counts.get('C',0), fracs.get('H_frac',0.0), fracs.get('E_frac',0.0), fracs.get('C_frac',0.0)])
                        rama = gq.get('ramachandran', {}); clashes = gq.get('clashes', [])
                        clash_rows.append([pid, rama.get('favored',0), rama.get('allowed',0), rama.get('outliers',0), len(clashes)])
                        for p in pk.get('pockets', []) if pk else []:
                            pocket_rows.append([pid, p.get('id') or p.get('index'), p.get('volume'), p.get('residue_count')])
                        if hb and hb.get('counts'):
                            for subtype, c in hb.get('counts', {}).items():
                                hbond_rows.append([pid, subtype, c, hb.get('fractions', {}).get(subtype,0.0)])
                    ss_io= _io.StringIO(); w=_csv.writer(ss_io); w.writerow(["PDB_ID","Helix_Count","Sheet_Count","Coil_Count","Helix_Frac","Sheet_Frac","Coil_Frac"]); [w.writerow(r) for r in ss_rows]
                    zf.writestr("data/aggregate_secondary_structure.csv", ss_io.getvalue())
                    gq_io= _io.StringIO(); w2=_csv.writer(gq_io); w2.writerow(["PDB_ID","Rama_Favored","Rama_Allowed","Rama_Outliers","Clash_Count"]); [w2.writerow(r) for r in clash_rows]
                    zf.writestr("data/aggregate_geometry_quality.csv", gq_io.getvalue())
                    if pocket_rows:
                        pk_io=_io.StringIO(); w3=_csv.writer(pk_io); w3.writerow(["PDB_ID","Pocket_ID","Volume","Residue_Count"]); [w3.writerow(r) for r in pocket_rows]
                        zf.writestr("data/aggregate_pockets.csv", pk_io.getvalue())
                    if hbond_rows:
                        hb_io=_io.StringIO(); w4=_csv.writer(hb_io); w4.writerow(["PDB_ID","Subtype","Count","Fraction"]); [w4.writerow(r) for r in hbond_rows]
                        zf.writestr("data/aggregate_hbond_subtypes.csv", hb_io.getvalue())
            except Exception as e:
                zf.writestr("data/ERROR_aggregate.txt", f"Failed aggregate summaries: {e}")
            readme = ("Protein Interaction Analysis - Complete Package\n\n" f"Structures: {', '.join(pdb_ids)}\n" "Contents:\n" "- report/Protein_Interaction_Report.pdf (now includes provenance, detailed statistics, residue network)\n" "- report/Protein_Interaction_Presentation.pptx\n" "- data/Interaction_Summary.xlsx\n" "- data/interaction_statistics.json (per-type metrics: counts, distance/angle stats, strengths)\n" "- data/residue_network_hubs.csv (top residue hubs in interaction network)\n" "- data/<PDB>_residue_profiles.csv (per-residue interaction aggregation, if computed)\n" "- data/<PDB>_interfaces.csv (chain-chain interaction interfaces, if computed)\n" "- data/<PDB>_outliers.csv (near-threshold interactions, if computed)\n" "- data/<PDB>_motifs.csv (experimental motif detections, if computed)\n" "- latex/report.tex (publication-ready LaTeX sections)\n" "- data/<PDB>_interactions.csv (flat interaction list)\n" "- data/<PDB>_summary.json (raw serialized analysis)\n")
            zf.writestr("README.txt", readme)
        return zip_buffer.getvalue()

    # ------------------------------------------------------------------
    # Internal Helpers
    # ------------------------------------------------------------------
    def _ensure_all_extensions(self,
                               pdb_ids: List[str],
                               analysis_results: Dict[str, Dict[str, Any]]):
        """Compute all enabled extensions for each structure if missing.

        This mirrors the lazy UI computation pathway (MainInterface._ensure_extensions)
        but keeps the logic local so packaging can be self-sufficient.
        """
        # Build mapping of (config flag attr, extension key, compute function)
        extension_specs = [
            ('enable_residue_profiles', 'residue_profiles', compute_residue_profiles),
            ('enable_interface_analysis', 'interface_analysis', compute_interface_analysis),
            ('enable_outlier_detection', 'outliers', compute_outliers),
            ('enable_provenance_panel', 'provenance', compute_provenance),
            ('enable_motif_detection', 'motifs', compute_motifs),
            # Structural / quality
            ('enable_secondary_structure', 'secondary_structure', compute_secondary_structure),
            ('enable_sasa_bsa', 'sasa_bsa', compute_sasa_bsa),
            ('enable_geometry_quality', 'geometry_quality', compute_geometry_quality),
            ('enable_disulfide_analysis', 'disulfides', compute_disulfides),
            ('enable_pocket_detection', 'pockets', compute_pockets),
            ('enable_conservation', 'conservation', compute_conservation),
            ('enable_pi_pi_refinement', 'pi_pi_refinement', compute_pi_pi_refinement),
            ('enable_hbond_subtypes', 'hbond_subtypes', compute_hbond_subtypes),
        ]

        for pdb_id in pdb_ids:
            result = analysis_results.get(pdb_id)
            if not result:
                continue
            exts = result.setdefault('extensions', {})
            for flag, key, func in extension_specs:
                if not getattr(self.config, flag, False):
                    continue
                if key in exts:  # already computed
                    continue
                if func is None:
                    # Module unavailable (optional dependency); skip gracefully
                    continue
                try:
                    exts[key] = func(result, self.config)
                except Exception as e:  # pragma: no cover - robustness
                    exts[key] = {'error': f'Failed to compute extension {key}: {e}'}
    
    def _add_structure_slide(self, prs, pdb_id: str, result: Dict[str, Any]):
        """Add slide for individual structure."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = f'Structure: {pdb_id}'
        
        tf = body_shape.text_frame
        interactions = result.get('interactions', {})
        total = sum(len(int_list) for int_list in interactions.values())
        
        tf.text = f'Total Interactions: {total}'
        
        for int_type, int_list in interactions.items():
            if int_list:
                p = tf.add_paragraph()
                display_name = self._get_interaction_display_name(int_type)
                p.text = f'{display_name}: {len(int_list)}'
    
    # Advanced LaTeX export (retained earlier implementation name kept for compatibility)
    def generate_latex_export(self,
                              pdb_ids: List[str],
                              analysis_results: Dict[str, Dict[str, Any]]) -> str:
        cfg = self.config.interactions
        stats = self._compute_interaction_statistics(analysis_results)
        latex_content: List[str] = []
        latex_content.append(r"\section{Protein Interaction Analysis}")
        latex_content.append("")
        latex_content.append(r"\subsection{Interaction Summary}")
        latex_content.append(r"\begin{table}[h]")
        latex_content.append(r"\centering")
        latex_content.append(r"\caption{Per-structure totals of detected noncovalent interactions}")
        latex_content.append(r"\begin{tabular}{|l|c|}")
        latex_content.append(r"\hline")
        latex_content.append(r"\textbf{Structure} & \textbf{Total Interactions} \\")
        latex_content.append(r"\hline")
        for pdb_id in pdb_ids:
            if pdb_id in analysis_results and analysis_results[pdb_id]:
                interactions = analysis_results[pdb_id].get('interactions', {})
                total = sum(len(int_list) for int_list in interactions.values())
                latex_content.append(f"{pdb_id} & {total} \\")
        latex_content.append(r"\hline")
        latex_content.append(r"\end{tabular}")
        latex_content.append(r"\label{tab:interaction_summary}")
        latex_content.append(r"\end{table}")
        latex_content.append("")
        if stats['interaction_types']:
            latex_content.append(r"\subsection{Detailed Interaction Statistics}")
            latex_content.append(r"\begin{table}[h]")
            latex_content.append(r"\centering")
            latex_content.append(r"\small")
            latex_content.append(r"\begin{tabular}{|l|r|r|r|r|r|r|r|}")
            latex_content.append(r"\hline")
            latex_content.append(r"Type & Total & Structures & Dist$_{mean}$ & Dist$_{std}$ & Ang$_{mean}$ & Strong & Moderate \\")
            latex_content.append(r"\hline")
            for int_type, rec in sorted(stats['interaction_types'].items(), key=lambda x: x[1]['total_count'], reverse=True):
                disp = self._get_interaction_display_name(int_type).replace('π', '$\\pi$')
                d_stats = rec['distance_stats'] or {}
                a_stats = rec['angle_stats'] or {}
                strengths = rec['strength_distribution']
                strong = strengths.get('strong', 0)
                moderate = strengths.get('moderate', 0)
                latex_content.append(
                    f"{disp} & {rec['total_count']} & {rec['structures_with']} & {d_stats.get('mean','-')} & {d_stats.get('std','-')} & {a_stats.get('mean','-')} & {strong} & {moderate} \\")
            latex_content.append(r"\hline")
            latex_content.append(r"\end{tabular}")
            latex_content.append(r"\label{tab:detailed_stats}")
            latex_content.append(r"\end{table}")
            latex_content.append("")
        latex_content.append(r"\subsection{Detection Parameters}")
        latex_content.append(r"\begin{itemize}")
        latex_content.append(fr"\item Hydrogen bonds: distance $\le$ {cfg.hbond_distance_cutoff}\AA, angle $\ge$ {cfg.hbond_angle_cutoff}$^\circ$")
        latex_content.append(fr"\item Halogen bonds: distance $\le$ {cfg.halogen_distance_cutoff}\AA, angle $\ge$ {cfg.halogen_angle_cutoff}$^\circ$")
        latex_content.append(fr"\item $\pi$-$\pi$ stacking: centroid distance $\le$ {cfg.pi_pi_distance_cutoff}\AA, plane angle $\le$ {cfg.pi_pi_angle_cutoff}$^\circ$")
        latex_content.append(fr"\item C-H$\cdots\pi$: distance window {cfg.ch_pi_min_distance}-{cfg.ch_pi_max_distance}\AA, angle $\le$ {cfg.ch_pi_angle_cutoff}$^\circ$")
        latex_content.append(fr"\item Anion-$\pi$: centroid distance $\le$ {cfg.anion_pi_distance_cutoff}\AA")
        latex_content.append(fr"\item n$\rightarrow\pi^*$: distance $\le$ {cfg.n_pi_star_distance_cutoff}\AA, angle $\ge$ {cfg.n_pi_star_angle_cutoff}$^\circ$")
        latex_content.append(fr"\item Ionic: charge center distance $\le$ {cfg.ionic_distance_cutoff}\AA")
        latex_content.append(fr"\item Hydrophobic: heavy atom distance $\le$ {cfg.hydrophobic_distance_cutoff}\AA (non-local)")
        latex_content.append(r"\end{itemize}")
        latex_content.append("")
        latex_content.append(r"\subsection{Methods}")
        latex_content.append("Noncovalent interactions were analyzed using the Protein Interaction Explorer ")
        latex_content.append(f"software (version {self.config.version}). Distance and angular thresholds follow literature-derived geometric criteria for each interaction class.")
        latex_content.append("")
        return "\n".join(latex_content)
    
    def _get_interaction_display_name(self, interaction_type: str) -> str:
        """Get human-readable name for interaction type."""
        display_names = {
            'hydrogen_bond': 'Hydrogen Bonds',
            'halogen_bond': 'Halogen Bonds',
            'pi_pi': 'π-π Stacking',
            'ionic': 'Ionic Interactions',
            'hydrophobic': 'Hydrophobic Contacts',
            'ch_pi': 'C-H···π Interactions',
            'chalcogen_bond': 'Chalcogen Bonds',
            'pnictogen_bond': 'Pnictogen Bonds',
            'tetrel_bond': 'Tetrel Bonds',
            'anion_pi': 'Anion-π Interactions',
            'n_pi_star': 'n→π* Interactions',
            'dispersion': 'London Dispersion'
        }
        return display_names.get(interaction_type, interaction_type.replace('_', ' ').title())
